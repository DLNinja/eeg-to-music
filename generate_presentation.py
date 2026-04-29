from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Music Generation Module"
subtitle.text = "Recent Modifications & Enhancements\nFixes #1 - #4"

# Fix 1
bullet_slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Fix #1: Emotions to Valence-Arousal Mapping"
tf = body_shape.text_frame
tf.text = "Mapped one main emotion to each corner of the Valence-Arousal (V-A) matrix."
p = tf.add_paragraph()
p.text = "This extreme positioning helps us use the most of our Markov chain model."

# Fix 2
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Fix #2: Markov Chain of Chords"
tf = body_shape.text_frame
tf.text = "Introduced a Markov chain of chords on relative scale degrees."
p = tf.add_paragraph()
p.text = "This is specifically tailored for each macro emotion."

# Fix 3
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Fix #3: Spike Profiles for Emotion Transitions"
tf = body_shape.text_frame
tf.text = "Added spike profiles for smooth emotion transitions (e.g., Happy -> Sad = Bittersweet)."
p = tf.add_paragraph()
p.text = "Each profile contains:"
p.level = 1
p = tf.add_paragraph()
p.text = "Transition emotion name"
p.level = 2
p = tf.add_paragraph()
p.text = "Tempo multiplier"
p.level = 2
p = tf.add_paragraph()
p.text = "Velocity shifts"
p.level = 2
p = tf.add_paragraph()
p.text = "Chord color"
p.level = 2
p = tf.add_paragraph()
p.text = "Melody register"
p.level = 2
p = tf.add_paragraph()
p.text = "Probability of a rest happening"
p.level = 2

# Fix 4
slide = prs.slides.add_slide(bullet_slide_layout)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Fix #4: Individual Emotion Fixes"
tf = body_shape.text_frame
tf.text = "Adjusted individual emotions for better macro mood music generation:"
p = tf.add_paragraph()
p.text = "Happy: Removed trills that repeat too often."
p.level = 1
p = tf.add_paragraph()
p.text = "Sadness: Solved musical conflicts between chords and melody."
p.level = 1
p = tf.add_paragraph()
p.text = "Fear: Introduced 3 distinct modes (Tension, Melodic, Climax)."
p.level = 1
p = tf.add_paragraph()
p.text = "Neutral: Placed in a different spot on the V-A matrix, giving it its own personality rather than trying to make it a musical chameleon."
p.level = 1

prs.save("Music_Generation_Module_Updates.pptx")
