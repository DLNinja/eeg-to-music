import sys
import subprocess
import os

print("Installing python-pptx if necessary...")
try:
    import pptx
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation

def main():
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "EEG2Music: Real-Time Emotional Synthesis"
    slide.placeholders[1].text = "The Generative Music Pipeline & Variable-Order Markov Logic"

    # Rest of the content
    slides_content = [
        ("1. System Architecture Overview", [
            "Emotion Tracker: Separates fast spikes from slow moods.",
            "Synthesizer: Builds chords, rhythms, and manages the melodic scale.",
            "Markov Engine: Provides musically realistic interval transitions.",
            "FluidSynth Audio: Generates the MIDI audio output in real-time."
        ]),
        ("2. The Emotion Tracker", [
            "Brainwave classifiers output 4 probabilities (Happy, Sad, Fear, Neutral) every second.",
            "Macro State (Mood): A slow, rolling exponential moving average of Valence and Arousal.",
            "  - Controls Key signature, scale mode, base tempo, default chord quality.",
            "Micro State (Spikes): The instantaneous deviation from the Macro State.",
            "  - Controls sudden volume boosts, chromatic melodic bends, sudden rhythmic density changes."
        ]),
        ("3. Dynamic Key Selection", [
            "Starting key maps dynamically to the first detected dominant emotion (Schubert-inspired).",
            "Happy: C Major / G Major (+0 / +7 offset)",
            "Sad: D Minor / F Minor (+2 / +5 offset)",
            "Fear: C# Minor / Eb Minor (+1 / +3 offset)",
            "Neutral: F Major / A Minor (+5 / +9 offset)",
            "Ensures tonal center aligns perfectly with the earliest tracked brain state."
        ]),
        ("4. Foundation: Rules & Templates", [
            "Only the Lead Melody uses stochastic Markov logic. The rest follows rigid rules for musical stability:",
            "Harmony (Chords): Driven by Macro Valence. Maps functional progressions (e.g. sad=[0,6,5,4]).",
            "Rhythmic Density: Driven by Macro Arousal. Selects strict subdivision templates.",
            "Dynamics: Driven by Micro Arousal spikes to create sudden swells in velocity."
        ]),
        ("5. The Challenge: Melodic Intent", [
            "Simple 1st-Order Markov Chains only remember the very last melodic interval played.",
            "The Problem: If the last jump was +2, the model doesn't know if that was part of a rising scale or falling sequence.",
            "The Result: The melody creates an endless 'random walk'. It feels short-sighted and lacks structural phrasing."
        ]),
        ("6. The Solution: N-Gram Models", [
            "Variable-Order Markov Chain implemented to solve aimlessness.",
            "Extracts intervals up to the 3rd order directly from human composers in VGMIDI.",
            "  - 1st Order: '+2'",
            "  - 2nd Order: '-1, +2'",
            "  - 3rd Order: '+4, -1, +2'",
            "Models exported as nested JSON matrices grouped by emotional quadrant."
        ]),
        ("7. Variable-Order Query Engine", [
            "The MarkovEngine uses a best-match fallback strategy.",
            "When Synthesizer asks for a note, it provides the last 3 intervals: [+4, -1, +2].",
            "1. Attempt 3rd Order: Does model know what follows [+4, -1, +2]? -> Long-term melodic intent.",
            "2. Attempt 2nd Order: If not, what follows [-1, +2]? -> Shorter intent.",
            "3. Attempt 1st Order: If not, what follows [+2]? -> Reliable safety net.",
            "4. Absolute Fallback: If Unseen, stand still (0)."
        ]),
        ("8. Handling Emotional Spikes", [
            "Variable-Order loops cause 'State Mismatches' during emotional shifts.",
            "If emotion shifts Happy -> Fear, querying the Fear matrix with Happy intervals results in dissonant out-of-scale leaps.",
            "The Fix applied to real-time generator:",
            "  - Wipes interval history back to [0, 0, 0] (Stationary) on mode shifts and spikes.",
            "  - Forces a Harmonic Snap: Next note bypasses Markov and snaps safely to the Root/3rd/5th of the new chord."
        ]),
        ("9. Real-Time Audio Synthesis", [
            "Sequences pushed to an asynchronous queue for 1-second EEG windows.",
            "Precise Timing: Background thread checks 2ms-resolution clock and triggers FluidSynth exactly when due.",
            "Lag Compensation: Audio playback starts with a stable 1-sec buffer delay to absorb processing jitter.",
            "Stability: FluidSynth configured to explicitly ignore MIDI inputs, preventing Windows driver crash loops."
        ])
    ]

    for title_text, bullets in slides_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title_text
        tf = slide.placeholders[1].text_frame
        tf.clear()  # Clear default empty paragraph
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            if p.text.startswith("  - "):
                p.level = 1
                p.text = p.text.replace("  - ", "")
            else:
                p.level = 0

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "EEG2Music_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    main()
